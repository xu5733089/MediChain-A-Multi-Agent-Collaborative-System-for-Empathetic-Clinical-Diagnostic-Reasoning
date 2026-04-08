"""
MediChain_DemoB.pptx generator
Mirrors MediChain_DemoB.html — 23 slides, dark navy theme, entrance animations.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree
import os

# ── Constants ──────────────────────────────────────────────────────────────────
W, H = Inches(13.33), Inches(7.5)   # 1920×1080 aspect

# Colors
BG      = RGBColor(0x08, 0x0F, 0x1A)
BG2     = RGBColor(0x0D, 0x1F, 0x35)
TEAL    = RGBColor(0x00, 0xE5, 0xC0)
TEAL2   = RGBColor(0x00, 0xC4, 0xA5)
TEXT    = RGBColor(0xE8, 0xF4, 0xF8)
TEXT2   = RGBColor(0xA8, 0xC8, 0xD8)
MUTED   = RGBColor(0x6A, 0x9B, 0xB5)
AMBER   = RGBColor(0xF5, 0xA6, 0x23)
ROSE    = RGBColor(0xFF, 0x6B, 0x6B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
NAVY    = RGBColor(0x16, 0x32, 0x68)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ───────────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(blank)
    # Dark background rectangle
    bg = sl.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    # Subtle grid overlay (very faint lines — drawn as thin rectangles)
    return sl

def txb(sl, text, x, y, w, h, size=18, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True, font="Calibri"):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb

def rect(sl, x, y, w, h, fill=BG2, line_color=None, line_w=Pt(0.5)):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def hline(sl, x, y, w, color=TEAL, thickness=Pt(1.5)):
    ln = sl.shapes.add_shape(1, x, y, w, Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln

def badge(sl, text, x, y, color=TEAL):
    w = Inches(len(text) * 0.09 + 0.4)
    bg_c = RGBColor(int(color[0]*0.15), int(color[1]*0.15), int(color[2]*0.15))
    r = rect(sl, x, y, w, Inches(0.28), fill=RGBColor(0x00, 0x2A, 0x22), line_color=color, line_w=Pt(0.75))
    txb(sl, text, x + Inches(0.08), y + Inches(0.04), w, Inches(0.22), size=9, color=color, font="Consolas")
    return r

def eyebrow(sl, text):
    txb(sl, f"[ {text} ]", Inches(0.6), Inches(0.22), Inches(8), Inches(0.35),
        size=11, color=TEAL, font="Consolas")

def slide_num(sl, n, total=23):
    txb(sl, f"{n:02d} / {total}", Inches(11.5), Inches(0.22), Inches(1.5), Inches(0.35),
        size=11, color=MUTED, align=PP_ALIGN.RIGHT, font="Consolas")

def title_line(sl, t1, t2="", y=Inches(0.7)):
    """Big serif-style title. t2 is the italic teal part."""
    txb(sl, t1, Inches(0.6), y, Inches(9), Inches(1.0), size=52, bold=False, color=TEXT, font="Georgia")
    if t2:
        txb(sl, t2, Inches(0.6 + len(t1)*0.29), y, Inches(9), Inches(1.0),
            size=52, bold=False, italic=True, color=TEAL, font="Georgia")

def subtitle(sl, text, y=Inches(1.75)):
    txb(sl, text, Inches(0.6), y, Inches(12), Inches(0.5),
        size=15, color=MUTED, italic=True, font="Calibri")

def bottom_bar(sl):
    hline(sl, 0, H - Inches(0.04), W, TEAL, Pt(3))

def card_box(sl, x, y, w, h, accent=TEAL):
    r = rect(sl, x, y, w, h, fill=BG2, line_color=RGBColor(0x00, 0x40, 0x35), line_w=Pt(0.5))
    # top accent line
    rect(sl, x, y, w, Pt(2), fill=accent)
    return r

def bullet_row(sl, icon, title, desc, x, y, w, icon_color=TEAL):
    txb(sl, "→", x, y, Inches(0.25), Inches(0.3), size=13, color=icon_color, font="Consolas")
    txb(sl, title, x + Inches(0.28), y, w - Inches(0.28), Inches(0.28), size=14, bold=True, color=TEXT2, font="Calibri")
    txb(sl, desc, x + Inches(0.28), y + Inches(0.26), w - Inches(0.28), Inches(0.32), size=12, color=MUTED, font="Calibri")

def pipe_step(sl, num, title, desc, x, y, w):
    card_box(sl, x, y, w, Inches(0.9), accent=TEAL)
    txb(sl, num, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.6),
        size=12, color=TEAL, font="Consolas")
    txb(sl, title, x + Inches(0.55), y + Inches(0.1), w - Inches(0.7), Inches(0.35),
        size=15, bold=True, color=TEXT, font="Calibri")
    txb(sl, desc, x + Inches(0.55), y + Inches(0.44), w - Inches(0.7), Inches(0.42),
        size=11, color=MUTED, font="Calibri", wrap=True)

def stat_card(sl, value, label, x, y, w=Inches(2.8), color=TEAL):
    card_box(sl, x, y, w, Inches(1.1), accent=color)
    txb(sl, value, x + Inches(0.15), y + Inches(0.1), w - Inches(0.2), Inches(0.6),
        size=38, bold=False, color=color, font="Georgia")
    txb(sl, label, x + Inches(0.15), y + Inches(0.68), w - Inches(0.2), Inches(0.35),
        size=10, color=MUTED, font="Consolas")

def add_animation(shape, effect="appear", delay_ms=0, duration_ms=500):
    """Add entrance animation to a shape via XML."""
    sp_id = shape.shape_id
    slide_xml = shape._element.getparent().getparent()

    # Find or create timing element
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    timing = slide_xml.find(f"{{{ns}}}timing")
    if timing is None:
        timing = etree.SubElement(slide_xml, f"{{{ns}}}timing")

    tnLst = timing.find(f"{{{ns}}}tnLst")
    if tnLst is None:
        tnLst = etree.SubElement(timing, f"{{{ns}}}tnLst")

    # Build animation XML
    delay_str = str(delay_ms) if delay_ms > 0 else "0"
    dur_str = str(duration_ms)

    anim_xml = f"""
    <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
           xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <p:cTn id="{sp_id * 10}" dur="indefinite" nodeType="clickEffect">
        <p:stCondLst>
          <p:cond delay="{delay_str}"/>
        </p:stCondLst>
        <p:childTnLst>
          <p:par>
            <p:cTn id="{sp_id * 10 + 1}" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
              <p:childTnLst>
                <p:par>
                  <p:cTn id="{sp_id * 10 + 2}" presetID="1" presetClass="entr" presetSubtype="0"
                         fill="hold" grpId="0" nodeType="clickEffect" dur="{dur_str}">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>
                      <p:set>
                        <p:cBhvr>
                          <p:cTn id="{sp_id * 10 + 3}" dur="1" fill="hold"/>
                          <p:tgtEl>
                            <p:spTgt spid="{sp_id}"/>
                          </p:tgtEl>
                          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:to><p:strVal val="visible"/></p:to>
                      </p:set>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:par>
        </p:childTnLst>
      </p:cTn>
    </p:par>
    """
    try:
        par_el = etree.fromstring(anim_xml)
        tnLst.append(par_el)
    except Exception:
        pass


# ── SLIDE 1: Title ─────────────────────────────────────────────────────────────
sl = add_slide()
# Glow orb (faint circle)
orb = sl.shapes.add_shape(9, Inches(8.5), Inches(-1), Inches(6), Inches(6))  # oval
orb.fill.solid(); orb.fill.fore_color.rgb = RGBColor(0x00, 0x1A, 0x28)
orb.line.fill.background()

txb(sl, "COMP9900 · Team 9900-W18C-CAKE · UNSW Sydney",
    Inches(0.6), Inches(1.0), Inches(9), Inches(0.4),
    size=12, color=TEAL, font="Consolas")

# Big title
txb(sl, "Medi", Inches(0.6), Inches(1.55), Inches(5.5), Inches(2.2),
    size=108, color=TEXT, font="Georgia")
txb(sl, "Chain", Inches(3.05), Inches(1.55), Inches(6), Inches(2.2),
    size=108, italic=True, color=TEAL, font="Georgia")

txb(sl, "A Multi-Agent Collaborative System for\nEmpathetic Clinical Diagnostic Reasoning",
    Inches(0.6), Inches(3.85), Inches(8.5), Inches(0.9),
    size=20, color=MUTED, italic=True, font="Calibri")

hline(sl, Inches(0.6), Inches(4.85), Inches(4.5), TEAL)

# Meta stats row
for i, (val, lbl) in enumerate([
    ("48 SP", "Story Points"), ("47,441", "Medical Docs"),
    ("6", "Modalities"), ("3", "Languages")
]):
    x = Inches(0.6 + i * 2.2)
    txb(sl, val, x, Inches(5.1), Inches(2.0), Inches(0.7),
        size=34, color=TEAL, font="Georgia")
    txb(sl, lbl, x, Inches(5.72), Inches(2.0), Inches(0.35),
        size=10, color=MUTED, font="Consolas")

# Rings graphic (simple concentric circles)
for r_size, opacity_hex in [(1.8, "1A"), (1.3, "14"), (0.8, "22"), (0.3, "FF")]:
    cx, cy = Inches(10.8), Inches(3.5)
    rs = Inches(r_size)
    sh = sl.shapes.add_shape(9, cx - rs, cy - rs * 0.7, rs * 2, rs * 1.4)
    sh.fill.background()
    sh.line.color.rgb = TEAL
    sh.line.width = Pt(0.75)

txb(sl, "Mc", Inches(10.45), Inches(3.2), Inches(0.7), Inches(0.5),
    size=22, color=BG, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
core = sl.shapes.add_shape(9, Inches(10.42), Inches(3.12), Inches(0.76), Inches(0.53))
core.fill.solid(); core.fill.fore_color.rgb = TEAL
core.line.fill.background()

txb(sl, "▶  FINAL PRESENTATION", Inches(0.6), Inches(0.55), Inches(4), Inches(0.38),
    size=12, bold=True, color=BG, font="Consolas")
badge_bg = rect(sl, Inches(0.6), Inches(0.55), Inches(2.8), Inches(0.38), fill=TEAL)

bottom_bar(sl)


# ── SLIDE 2: Problem Background ────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Motivation")
slide_num(sl, 2)
title_line(sl, "Problem ", "Background")
subtitle(sl, "Diagnostic errors remain one of healthcare's most persistent and costly failures")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

# 4 stat cards
stats = [
    ("12M+", "Diagnostic Errors/Year USA", "Graber et al., 2013", ROSE),
    ("22%",  "Error Rate · Internal Med.", "BMJ Quality & Safety, 2016", AMBER),
    ("$40B", "Annual Cost · Dx Failures", "CRICO Strategies, 2021", ROSE),
    ("71%",  "GPT-4 on MedQA bench.", "Kung et al., 2023", MUTED),
]
for i, (val, lbl, cite, col) in enumerate(stats):
    x = Inches(0.5 + i * 3.2)
    card_box(sl, x, Inches(2.4), Inches(3.0), Inches(2.2), accent=col)
    txb(sl, val, x + Inches(0.15), Inches(2.5), Inches(2.7), Inches(0.9),
        size=52, color=col, font="Georgia")
    txb(sl, lbl, x + Inches(0.15), Inches(3.35), Inches(2.7), Inches(0.4),
        size=13, bold=True, color=TEXT2, font="Calibri")
    txb(sl, cite, x + Inches(0.15), Inches(3.72), Inches(2.7), Inches(0.3),
        size=9, color=MUTED, font="Consolas")

# Gap box
rect(sl, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.5),
     fill=RGBColor(0x1A, 0x0E, 0x00), line_color=AMBER, line_w=Pt(0.75))
txb(sl, "Critical Gap:",
    Inches(0.7), Inches(4.9), Inches(2.2), Inches(0.4), size=14, bold=True, color=AMBER)
txb(sl, "GPT-4's 71% MedQA accuracy falls below the >90% clinical threshold. "
        "On rare disease classification the gap exceeds 50 percentage points.",
    Inches(0.7), Inches(5.28), Inches(12.0), Inches(0.9), size=13, color=MUTED, wrap=True)

bottom_bar(sl)


# ── SLIDE 3: Existing Limitations ──────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Research Gap")
slide_num(sl, 3)
title_line(sl, "Existing ", "Limitations")
subtitle(sl, "Current AI diagnostic tools fail on at least one of five critical clinical dimensions")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

# Table
headers = ["Capability", "GPT-4 Standalone", "RAG-only", "Single-Agent", "MediChain (ours)"]
col_colors = [TEXT2, ROSE, MUTED, AMBER, TEAL]
col_w = [Inches(2.8), Inches(2.2), Inches(2.0), Inches(2.0), Inches(2.5)]
rows = [
    ["Hallucination Control", "✗ None", "△ Partial", "△ Partial", "✓ RAG + Critic"],
    ["Explainable Reasoning", "✗ Black-box", "✗ None", "△ Single-pass", "✓ Agent chain"],
    ["Clinical Safety Rules", "✗ None", "✗ None", "✗ None", "✓ Cardiac·Stroke·Drug"],
    ["Multimodal Intake",     "△ Image only", "✗ Text only", "✗ Text only", "✓ 6 modalities"],
    ["Empathetic Dialogue",   "✗ None", "✗ No dialogue", "✗ None", "✓ SOCRATES + warmth"],
    ["Knowledge Currency",    "✗ Static cutoff", "✓ Dynamic", "✗ Static", "✓ 47,441 docs live"],
    ["Multi-Disciplinary",    "✗ Single pass", "✗ No review", "✗ Single", "✓ 3 agents"],
    ["Verified Citations",    "✗ None", "△ Retrieval", "✗ None", "✓ PubMed per Dx"],
]
ys = 2.35
x_start = 0.5
for ci, (hdr, col) in enumerate(zip(headers, col_colors)):
    x = Inches(x_start + sum(w.inches for w in col_w[:ci]))
    txb(sl, hdr, x, Inches(ys), col_w[ci], Inches(0.38),
        size=11, bold=True, color=col, font="Consolas")

hline(sl, Inches(0.5), Inches(2.72), Inches(11.8), TEAL, Pt(0.75))
for ri, row in enumerate(rows):
    y = Inches(2.8 + ri * 0.52)
    bg_color = RGBColor(0x0A, 0x18, 0x28) if ri % 2 == 0 else BG
    rect(sl, Inches(0.5), y, Inches(12.3), Inches(0.5), fill=bg_color)
    for ci, (cell, col_w_i) in enumerate(zip(row, col_w)):
        x = Inches(x_start + sum(w.inches for w in col_w[:ci]))
        is_good = cell.startswith("✓")
        is_bad  = cell.startswith("✗")
        col = TEAL if is_good else (ROSE if is_bad else AMBER)
        txb(sl, cell, x + Inches(0.05), y + Inches(0.08), col_w_i - Inches(0.1), Inches(0.38),
            size=11, color=col if ci > 0 else TEXT2, font="Calibri")

bottom_bar(sl)


# ── SLIDE 4: Design Pillars ─────────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Our Approach")
slide_num(sl, 4)
title_line(sl, "MediChain ", "Design")
subtitle(sl, "Four integrated design pillars address each identified limitation simultaneously")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

pillars = [
    ("🤝", "Multi-Agent Collaboration", "Clinical Team Simulation", TEAL,
     ["Interviewer → SOCRATES structured intake",
      "Diagnostician → ranked differential + evidence",
      "Critic → safety verification + hallucination check"]),
    ("📚", "Knowledge Grounding (RAG)", "Evidence-Based Reasoning", AMBER,
     ["47,441 PubMed abstracts in Qdrant vector DB",
      "BioLORD-2023 768d + BM25 hybrid search",
      "PubMed citation with relevance score per claim"]),
    ("🔍", "Explainable Reasoning", "Clinical Auditability", MUTED,
     ["Full agent reasoning chain in Results tab",
      "Critic verdict alongside every diagnosis",
      "Session exportable as PDF / JSON"]),
    ("💙", "Empathetic Interaction", "Patient-Centred Design", ROSE,
     ["Warm language — acknowledges before questioning",
      "LLM rewrites lay language to MeSH terms",
      "Multilingual: EN / 中文 / 日本語"]),
]

for i, (icon, name, role, col, feats) in enumerate(pillars):
    col_i = i % 2
    row_i = i // 2
    x = Inches(0.5 + col_i * 6.4)
    y = Inches(2.45 + row_i * 2.35)
    card_box(sl, x, y, Inches(6.1), Inches(2.18), accent=col)
    txb(sl, icon, x + Inches(0.15), y + Inches(0.12), Inches(0.5), Inches(0.5), size=26)
    txb(sl, name, x + Inches(0.7),  y + Inches(0.1),  Inches(5.2), Inches(0.38),
        size=18, bold=True, color=TEXT, font="Calibri")
    txb(sl, role, x + Inches(0.7),  y + Inches(0.45), Inches(5.2), Inches(0.28),
        size=10, color=col, font="Consolas")
    for j, feat in enumerate(feats):
        txb(sl, f"→  {feat}", x + Inches(0.15), y + Inches(0.85 + j * 0.4),
            Inches(5.8), Inches(0.36), size=12, color=MUTED)

bottom_bar(sl)


# ── SLIDE 5: System Architecture ───────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "System Design")
slide_num(sl, 5)
title_line(sl, "System ", "Architecture")
subtitle(sl, "Five-layer architecture — each layer independently testable and horizontally scalable")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

layers = [
    ("User",      [("React + Vite Frontend", "Symptom Input · Chat UI · Results · History", True),
                   ("Multimodal Zone",        "DICOM · Video · OCR · Audio · PDF · Compare", True),
                   ("i18n + RBAC",            "EN / ZH / JA · Patient view · Provider view", True)]),
    ("API",       [("FastAPI REST",           "Session routing · File upload · Export", False),
                   ("Multimodal API",         "/analyze/file · /analyze/ocr · /analyze/compare", True),
                   ("JWT Auth",               "Role claims · /api/provider/* guard", True)]),
    ("Agents",    [("Interviewer Agent",      "SOCRATES · Empathetic framing · Multimodal context", False),
                   ("Diagnostician Agent",    "Differential Dx · RAG citations · Confidence ranking", False),
                   ("Critic Agent",           "Safety rules · Hallucination check · Drug conflicts", True)]),
    ("Knowledge", [("Qdrant Vector DB",       "47,441 docs · BioLORD 768d Dense + BM25 Sparse + RRF", True),
                   ("LLM Rewriter",           "Patient lang → SNOMED/MeSH terms · 3-variant expansion", True),
                   ("SQLite Persistence",     "Sessions · Messages · Users · Uploads", False)]),
]

for li, (lbl, cards) in enumerate(layers):
    y = Inches(2.4 + li * 1.1)
    txb(sl, lbl, Inches(0.1), y + Inches(0.3), Inches(0.8), Inches(0.5),
        size=9, color=MUTED, font="Consolas")
    for ci, (name, desc, is_new) in enumerate(cards):
        x = Inches(1.0 + ci * 4.1)
        col = TEAL if is_new else RGBColor(0x1A, 0x35, 0x52)
        card_box(sl, x, y, Inches(3.9), Inches(0.95), accent=col)
        txb(sl, name + (" ★" if is_new else ""),
            x + Inches(0.12), y + Inches(0.1), Inches(3.65), Inches(0.38),
            size=13, bold=True, color=TEXT, font="Calibri")
        txb(sl, desc, x + Inches(0.12), y + Inches(0.46), Inches(3.65), Inches(0.42),
            size=10, color=MUTED, font="Calibri", wrap=True)
    if li < len(layers) - 1:
        txb(sl, "▼", Inches(6.3), y + Inches(0.95), Inches(0.5), Inches(0.3),
            size=14, color=MUTED, align=PP_ALIGN.CENTER)

bottom_bar(sl)


# ── SLIDE 6: Three-Agent Design ────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Core Innovation")
slide_num(sl, 6)
title_line(sl, "Three-Agent ", "Design")
subtitle(sl, "Specialised agents operate in sequence, mirroring a real multi-disciplinary clinical team")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

agents = [
    ("🩺", "Interviewer", "Patient Intake Specialist", TEAL,
     ["Opens with empathetic acknowledgement",
      "SOCRATES: Site · Onset · Character · Radiation",
      "Accepts text, voice, video, images as context",
      "Decides when sufficient info collected"],
     "Claude claude-sonnet-4-6"),
    ("🔬", "Diagnostician", "Evidence-Based Diagnosis", AMBER,
     ["Receives full symptom context from Interviewer",
      "Fires 3–4 Qdrant hybrid queries via LLM rewrite",
      "Produces ranked differential: condition · prob. · evidence",
      "Attaches PubMed citation for every claim"],
     "RAG-augmented · 47K docs"),
    ("⚖️", "Critic", "Clinical Safety Verifier", ROSE,
     ["Reviews Diagnostician's full reasoning chain",
      "Checks: logic · evidence alignment · hallucination",
      "Hard-coded triggers: STEMI · Stroke FAST · Drug Rx",
      "Issues: 'Consistent' or specific objection"],
     "Hard-coded safety rules"),
]

for i, (icon, name, role, col, feats, badge_txt) in enumerate(agents):
    x = Inches(0.5 + i * 4.3)
    card_box(sl, x, Inches(2.4), Inches(4.1), Inches(4.6), accent=col)
    # top border
    rect(sl, x, Inches(2.4), Inches(4.1), Inches(0.07), fill=col)
    txb(sl, icon, x + Inches(0.2), Inches(2.5), Inches(0.6), Inches(0.6), size=30)
    txb(sl, name, x + Inches(0.2), Inches(3.05), Inches(3.8), Inches(0.5),
        size=24, color=TEXT, font="Georgia")
    txb(sl, role, x + Inches(0.2), Inches(3.5), Inches(3.8), Inches(0.3),
        size=10, color=col, font="Consolas")
    hline(sl, x + Inches(0.2), Inches(3.82), Inches(3.7), col, Pt(0.5))
    for j, feat in enumerate(feats):
        txb(sl, f"→  {feat}", x + Inches(0.2), Inches(3.95 + j * 0.4),
            Inches(3.8), Inches(0.36), size=12, color=MUTED)
    badge(sl, badge_txt, x + Inches(0.2), Inches(6.6), col)

# Flow line
txb(sl, "Patient  →  Interviewer  →  Diagnostician  →  Critic  →  Results + Citations",
    Inches(0.5), Inches(7.1), Inches(12), Inches(0.35),
    size=12, color=MUTED, font="Consolas", align=PP_ALIGN.CENTER)

bottom_bar(sl)


# ── SLIDE 7: Delivery Summary ───────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Project Delivery")
slide_num(sl, 7)
title_line(sl, "Delivery ", "Summary")
subtitle(sl, "14 user stories · 48 story points · 2 sprints · all targets met on schedule")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

rect(sl, Inches(0.5), Inches(2.38), Inches(6.1), Inches(0.38),
     fill=RGBColor(0x00, 0x14, 0x10), line_color=TEAL, line_w=Pt(0.5))
txb(sl, "SPRINT 1 — Weeks 3–5 · 13 SP", Inches(0.65), Inches(2.42), Inches(5.8), Inches(0.3),
    size=12, color=TEAL, font="Consolas")
rect(sl, Inches(6.8), Inches(2.38), Inches(6.0), Inches(0.38),
     fill=RGBColor(0x18, 0x10, 0x00), line_color=AMBER, line_w=Pt(0.5))
txb(sl, "SPRINT 2 — Weeks 5–8 · 35 SP", Inches(6.95), Inches(2.42), Inches(5.7), Inches(0.3),
    size=12, color=AMBER, font="Consolas")

stories = [
    ("PROJ-1",    "Symptom Input UI + Session Init API", "4 SP", TEAL),
    ("PROJ-2–3",  "Patient Chat UI + Multi-Agent Coordinator Logic", "5 SP", TEAL),
    ("PROJ-4–6",  "Reasoning Display Panel + DB Schema + RAG Foundation", "4 SP", TEAL),
    ("PROJ-7",    "RAG 2.0 — Qdrant Hybrid Search (Dense + BM25 + RRF)", "8 SP", AMBER),
    ("PROJ-8",    "Multimodal Input — DICOM, Audio, Video, PDF, OCR", "8 SP", AMBER),
    ("PROJ-9",    "Dual Authentication (Patient / Provider) + JWT RBAC", "5 SP", AMBER),
    ("PROJ-10",   "AI Image Annotation — 9-region JSON overlay", "4 SP", AMBER),
    ("PROJ-11–14","Safety Rules · Empathetic AI · i18n (3 langs) · Session Export", "10 SP", AMBER),
]

for i, (pid, name, sp, col) in enumerate(stories):
    y = Inches(2.88 + i * 0.45)
    rect(sl, Inches(0.5), y, Inches(12.3), Inches(0.42),
         fill=RGBColor(0x0A, 0x18, 0x28) if i%2==0 else BG2,
         line_color=col, line_w=Pt(0.3))
    txb(sl, pid,  Inches(0.6),  y + Inches(0.07), Inches(1.2), Inches(0.3), size=10, color=col, font="Consolas")
    txb(sl, name, Inches(1.85), y + Inches(0.07), Inches(9.0), Inches(0.3), size=12, color=TEXT2, font="Calibri")
    txb(sl, sp,   Inches(11.1), y + Inches(0.07), Inches(0.9), Inches(0.3), size=10, color=MUTED, font="Consolas")
    txb(sl, "✓",  Inches(12.1), y + Inches(0.07), Inches(0.5), Inches(0.3), size=12, color=TEAL)

# Summary bar
rect(sl, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.72),
     fill=RGBColor(0x00, 0x14, 0x10), line_color=TEAL, line_w=Pt(0.75))
for i, (val, lbl) in enumerate([("48","Total SP"),("14","User Stories"),
                                  ("47,441","RAG Docs"),("6","Modalities"),
                                  ("3","Languages"),("4","New Endpoints"),("2","Auth Roles")]):
    x = Inches(0.8 + i * 1.73)
    txb(sl, val, x, Inches(6.58), Inches(1.5), Inches(0.38), size=22, color=TEAL, font="Georgia")
    txb(sl, lbl, x, Inches(6.93), Inches(1.5), Inches(0.28), size=9, color=MUTED, font="Consolas")

bottom_bar(sl)


# ── SLIDE 8: RAG 2.0 ───────────────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Feature · PROJ-7")
slide_num(sl, 8)
title_line(sl, "RAG ", "2.0")
subtitle(sl, "ChromaDB → Qdrant · BioLORD-2023 768d + BM25 + RRF · 9.4× knowledge growth")
hline(sl, Inches(0.6), Inches(2.15), Inches(2))

# Stat grid
s4 = [("9.4×","Knowledge growth"),("2×","Richer embeddings"),
      ("3","Search methods fused"),("4","Query variants/search")]
for i, (v, l) in enumerate(s4):
    stat_card(sl, v, l, Inches(0.5 + i*3.2), Inches(2.4), Inches(3.0), TEAL)

# Comparison bars (horizontal)
comparisons = [
    ("Knowledge Base", "~5K", "47,441", 0.105),
    ("Embedding Dims", "384d", "768d", 0.5),
    ("Query Variants", "1", "3–4", 0.25),
    ("Search Methods", "Cosine only", "Dense+BM25+RRF", 0.33),
    ("Recall Coverage", "thr=0.25", "thr=0.15 ↑", 0.6),
]
card_box(sl, Inches(0.5), Inches(3.72), Inches(8.0), Inches(3.35), accent=TEAL)
txb(sl, "Demo A vs Demo B — Key Metrics",
    Inches(0.65), Inches(3.8), Inches(7.5), Inches(0.35), size=11, color=MUTED, font="Consolas")
for i, (lbl, va, vb, ratio) in enumerate(comparisons):
    y = Inches(4.25 + i * 0.55)
    txb(sl, lbl, Inches(0.65), y, Inches(2.0), Inches(0.3), size=12, color=TEXT2, font="Calibri")
    # Demo A bar (red)
    bar_w = Inches(4.0 * ratio)
    rect(sl, Inches(2.7), y + Inches(0.02), bar_w, Inches(0.22),
         fill=RGBColor(0x80, 0x30, 0x30))
    txb(sl, va, Inches(2.7) + bar_w + Inches(0.05), y, Inches(1.5), Inches(0.28), size=10, color=ROSE, font="Consolas")
    # Demo B bar (teal) full
    rect(sl, Inches(2.7), y + Inches(0.26), Inches(4.0), Inches(0.22), fill=TEAL2)
    txb(sl, vb, Inches(6.75), y + Inches(0.24), Inches(1.8), Inches(0.28), size=10, color=TEAL, font="Consolas")

# Why boxes
why_items = [
    ("Why BioLORD-2023?",
     "Pre-trained on SNOMED CT, MeSH, UMLS. Natively understands \"heart attack\" = \"AMI\" = \"myocardial infarction\" without keyword overlap.", TEAL),
    ("Why Hybrid Search?",
     "Dense vectors match semantics; BM25 matches exact terms. RRF score = Σ 1/(k+rank) fuses both — outperforms either alone on biomedical IR.", TEAL),
    ("New: /api/rag/ingest  ★",
     "POST endpoint triggers on-demand PubMed fetch. 60+ ICD-10 search terms. Custom terms via request body. No restart required.", AMBER),
]
for i, (title, desc, col) in enumerate(why_items):
    y = Inches(3.72 + i * 1.1)
    card_box(sl, Inches(8.7), y, Inches(4.45), Inches(1.05), accent=col)
    txb(sl, title, Inches(8.85), y + Inches(0.08), Inches(4.1), Inches(0.3),
        size=13, bold=True, color=col, font="Calibri")
    txb(sl, desc,  Inches(8.85), y + Inches(0.38), Inches(4.1), Inches(0.6),
        size=11, color=MUTED, font="Calibri", wrap=True)

bottom_bar(sl)


# ── SLIDE 9: Query Pipeline ────────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "RAG Detail")
slide_num(sl, 9)
title_line(sl, "Query ", "Pipeline")
subtitle(sl, "6-stage pipeline — from patient language to PubMed-cited output, including multimodal image rewriting")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

steps = [
    ("00", "Image Findings → Medical Terms  ★",
     "Claude Vision → raw findings → rewrite_image_findings_for_rag() → SNOMED terms injected into query"),
    ("01", "LLM Query Rewrite",
     "\"chest tightens when climbing stairs\" → \"stable angina, exertional chest pain, cardiac ischaemia\""),
    ("02", "3-Variant Query Expansion",
     "① Raw text  ② Medical rewrite  ③ \"treatment and diagnosis of {term}\"  ④ \"What is {term}?\""),
    ("03", "Qdrant Hybrid Search",
     "Dense (BioLORD cosine) + BM25 sparse in parallel — semantic similarity AND exact clinical term matching"),
    ("04", "RRF Fusion → Top 6",
     "score = Σ 1/(k+rank) · Threshold 0.15 · Dedup by doc ID"),
    ("05", "Cited Diagnostic Output",
     "Condition · probability · PubMed ID · authors · year · relevance score · clickable link"),
]
for i, (num, title, desc) in enumerate(steps):
    pipe_step(sl, num, title, desc, Inches(0.5), Inches(2.38 + i * 0.83), Inches(7.8))

# Corpus breakdown
card_box(sl, Inches(8.5), Inches(2.38), Inches(4.65), Inches(4.95), accent=TEAL)
txb(sl, "Knowledge Corpus — 47,441 Docs",
    Inches(8.65), Inches(2.48), Inches(4.3), Inches(0.35), size=11, color=MUTED, font="Consolas")
corpus = [("PubMed Abstracts", "~41,000", 0.87, TEAL),
          ("Clinical Guidelines", "~5,000", 0.28, AMBER),
          ("Case Reports", "~1,441", 0.08, ROSE)]
for i, (lbl, cnt, ratio, col) in enumerate(corpus):
    y = Inches(2.95 + i * 0.85)
    txb(sl, lbl, Inches(8.65), y, Inches(4.2), Inches(0.28), size=12, color=TEXT2, font="Calibri")
    rect(sl, Inches(8.65), y + Inches(0.3), Inches(4.0 * ratio), Inches(0.28), fill=col)
    txb(sl, cnt, Inches(8.65 + 4.0 * ratio + 0.1), y + Inches(0.28),
        Inches(1.5), Inches(0.28), size=11, color=col, font="Consolas")

stat_card(sl, "768d", "Embedding dims", Inches(8.5), Inches(5.6), Inches(2.2), TEAL)
stat_card(sl, "36M",  "Float32 values", Inches(10.95), Inches(5.6), Inches(2.2), MUTED)

bottom_bar(sl)


# ── SLIDE 10: Multimodal Input ─────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Feature · PROJ-8")
slide_num(sl, 10)
title_line(sl, "Multimodal ", "Input")
subtitle(sl, "6 clinical input modalities — unified preprocessing pipeline → Claude Vision / Speech / Text agents")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

modalities = [
    ("🩻", "DICOM Medical Imaging", ".dcm → pydicom → pixel normalise → JPEG → Claude Vision · HU metadata"),
    ("🖼",  "Image Analysis + Annotation", "Claude Vision → ANNOTATIONS_JSON · 9-region grid · numbered dots"),
    ("🎙",  "Audio Transcription", "Upload MP3/WAV → SpeechRecognition → text → RAG"),
    ("📹", "Video Frame Consultation", "Webcam → canvas.toBlob() every 5 s → Claude Vision per frame"),
    ("📄", "PDF / Text Extraction", "pdfplumber → structured text → semantic embedding → Qdrant"),
    ("✍",  "OCR Handwriting", "Image upload → Claude Vision → structured clinical text"),
]
for i, (icon, title, desc) in enumerate(modalities):
    col_i = i % 3
    row_i = i // 3
    x = Inches(0.5 + col_i * 4.3)
    y = Inches(2.4 + row_i * 2.4)
    card_box(sl, x, y, Inches(4.1), Inches(2.25), accent=TEAL)
    txb(sl, icon, x + Inches(0.15), y + Inches(0.12), Inches(0.6), Inches(0.55), size=28)
    txb(sl, title, x + Inches(0.15), y + Inches(0.7),  Inches(3.85), Inches(0.4),
        size=14, bold=True, color=TEXT, font="Calibri")
    txb(sl, desc,  x + Inches(0.15), y + Inches(1.1),  Inches(3.85), Inches(0.9),
        size=11, color=MUTED, font="Calibri", wrap=True)
    badge(sl, "NEW", x + Inches(3.5), y + Inches(0.12), TEAL)

bottom_bar(sl)


# ── SLIDE 11: AI Image Annotation ─────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Feature · PROJ-10")
slide_num(sl, 11)
title_line(sl, "AI Image ", "Annotation")
subtitle(sl, "Claude Vision → structured JSON · 9-region grid · interactive numbered overlays")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

stat_card(sl, "9", "Annotation region zones (3×3 grid)", Inches(0.5), Inches(2.4), Inches(3.2), TEAL)

# 3×3 region grid
regions = ["UPPER LEFT","UPPER CTR","UPPER RIGHT",
           "MID LEFT","MID CTR","MID RIGHT",
           "LOWER LEFT","LOWER CTR","LOWER RIGHT"]
active_r = {0, 5, 6}
for i, reg in enumerate(regions):
    col_i = i % 3
    row_i = i // 3
    x = Inches(0.5 + col_i * 1.1)
    y = Inches(3.75 + row_i * 0.62)
    col = TEAL if i in active_r else RGBColor(0x1A, 0x35, 0x52)
    rect(sl, x, y, Inches(1.0), Inches(0.55),
         fill=RGBColor(0x00, 0x28, 0x22) if i in active_r else BG2,
         line_color=col, line_w=Pt(0.75))
    txb(sl, reg, x + Inches(0.04), y + Inches(0.12), Inches(0.92), Inches(0.3),
        size=8, color=col, align=PP_ALIGN.CENTER, font="Consolas")

# Annotation findings
findings = [
    ("①", ROSE,  "Upper-Left — Patchy opacity, consolidation (pneumonia vs. atelectasis)"),
    ("②", TEAL,  "Mid-Right pleura — Focal pleural thickening, monitor for effusion"),
    ("③", AMBER, "Lower-Left costophrenic — Blunting, lateral decubitus view suggested"),
]
for i, (num, col, desc) in enumerate(findings):
    y = Inches(5.9 + i * 0.48)
    txb(sl, num,  Inches(0.5), y, Inches(0.3), Inches(0.38), size=14, color=col, bold=True)
    txb(sl, desc, Inches(0.85), y, Inches(2.9), Inches(0.38), size=11, color=MUTED)

# Right column: technical detail
card_box(sl, Inches(4.0), Inches(2.4), Inches(9.0), Inches(2.3), accent=TEAL)
txb(sl, "Structured Output Format:", Inches(4.15), Inches(2.5), Inches(8.7), Inches(0.3),
    size=13, bold=True, color=TEAL, font="Calibri")
txb(sl, 'ANNOTATIONS_JSON:[{"region":"UPPER-LEFT","finding":"patchy opacity","severity":"moderate"},...]',
    Inches(4.15), Inches(2.85), Inches(8.7), Inches(0.5), size=11, color=TEAL2, font="Consolas")
txb(sl, "Parsed with json.loads() · validated against 9-region enum · dot markers rendered with CSS pulse animation · hover on dot highlights matching finding in list",
    Inches(4.15), Inches(3.35), Inches(8.7), Inches(0.7), size=12, color=MUTED, wrap=True)

card_box(sl, Inches(4.0), Inches(4.85), Inches(9.0), Inches(0.85), accent=AMBER)
txb(sl, "vs. Demo A:", Inches(4.15), Inches(4.95), Inches(1.5), Inches(0.28), size=13, bold=True, color=AMBER)
txb(sl, "Raw image displayed only, no overlay, no findings list. Demo B adds full structured annotation pipeline with Claude Vision JSON output.",
    Inches(5.7), Inches(4.95), Inches(7.1), Inches(0.65), size=12, color=MUTED, wrap=True)

for i, (lbl, col) in enumerate([("Multi-image compare",TEAL),("DICOM native",TEAL),
                                   ("Hover tooltips",AMBER),("JSON validated",TEAL)]):
    badge(sl, lbl, Inches(4.0 + i * 2.25), Inches(5.85), col)

bottom_bar(sl)


# ── SLIDE 12: Video Consultation ───────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Feature · Video")
slide_num(sl, 12)
title_line(sl, "Video Frame ", "Consultation")
subtitle(sl, "Real-time webcam consultation — simultaneous speech transcription and frame-by-frame visual analysis")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

vsteps = [
    ("01","getUserMedia() Camera Access",
     "Browser requests webcam permission · video stream rendered in modal · Web Speech API handles audio separately"),
    ("02","Continuous Speech Transcription",
     "Web Speech API — continuous + interimResults · lang from navigator.language · interim text overlaid on video"),
    ("03","Automatic Frame Capture",
     "canvas.toBlob() captures JPEG every 5 s (FRAME_INTERVAL_MS=5000) · POST /api/analyze/file · Claude Vision analyses each frame"),
    ("04","Stop → Bundle Context",
     "Transcript + all frame analyses combined → structured message → AI Interviewer as multimodal context"),
]
for i, (num, title, desc) in enumerate(vsteps):
    pipe_step(sl, num, title, desc, Inches(0.5), Inches(2.4 + i * 1.0), Inches(7.5))

for i, (v, l, col) in enumerate([
    ("5 s","Frame capture interval",TEAL), ("6","Max frames/session",TEAL),
    ("3","Speech lang options",AMBER), ("0","Extra npm deps",MUTED)
]):
    stat_card(sl, v, l, Inches(8.4 + (i%2)*2.35), Inches(2.4 + (i//2)*1.3), Inches(2.2), col)

card_box(sl, Inches(8.4), Inches(5.08), Inches(4.7), Inches(1.0), accent=TEAL)
txb(sl, "Clinical Value:", Inches(8.55), Inches(5.15), Inches(4.4), Inches(0.28), size=13, bold=True, color=TEAL)
txb(sl, "Patients with mobility limitations can describe symptoms verbally while showing affected areas on camera. Richer context than text alone.",
    Inches(8.55), Inches(5.44), Inches(4.4), Inches(0.55), size=11, color=MUTED, wrap=True)

card_box(sl, Inches(8.4), Inches(6.18), Inches(4.7), Inches(0.9), accent=AMBER)
txb(sl, "File-Only Submission:", Inches(8.55), Inches(6.25), Inches(4.4), Inches(0.28), size=13, bold=True, color=AMBER)
txb(sl, "Video provided → typed description optional. System auto-generates description from frame analyses on submit.",
    Inches(8.55), Inches(6.54), Inches(4.4), Inches(0.45), size=11, color=MUTED, wrap=True)

bottom_bar(sl)


# ── SLIDE 13: Empathetic AI ─────────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Feature · PROJ-14")
slide_num(sl, 13)
title_line(sl, "Empathetic ", "AI Interviewer")
subtitle(sl, "SOCRATES clinical framework + warm language conventions + LLM query rewriting for RAG")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

# Before/After cards
for is_after, (lbl, text, col) in enumerate([
    ("Generic (before):",
     '"What is the site of your pain? Rate severity 1–10. How long has it lasted?"',
     ROSE),
    ("Empathetic (after):",
     '"I\'m sorry you\'re dealing with this. Could you help me understand where the discomfort is strongest? Does it spread anywhere — like your arm, jaw, or back?"',
     TEAL),
]):
    y = Inches(2.45 + is_after * 1.25)
    card_box(sl, Inches(0.5), y, Inches(6.1), Inches(1.15), accent=col)
    txb(sl, lbl, Inches(0.65), y + Inches(0.1), Inches(5.8), Inches(0.28), size=11, color=col, font="Consolas")
    txb(sl, text, Inches(0.65), y + Inches(0.38), Inches(5.8), Inches(0.65), size=12, color=TEXT2, italic=True, wrap=True)

# RAG rewrite example
card_box(sl, Inches(0.5), Inches(5.05), Inches(6.1), Inches(1.0), accent=AMBER)
txb(sl, "RAG Query Rewrite:", Inches(0.65), Inches(5.12), Inches(5.8), Inches(0.28), size=11, color=AMBER, font="Consolas")
txb(sl, '"my chest hurts and i feel dizzy when i walk up stairs"',
    Inches(0.65), Inches(5.4), Inches(5.8), Inches(0.28), size=12, color=ROSE, italic=True)
txb(sl, '→  "stable angina, exertional chest pain, cardiac ischaemia, coronary artery disease"',
    Inches(0.65), Inches(5.68), Inches(5.8), Inches(0.3), size=11, color=TEAL, font="Consolas")

# SOCRATES
card_box(sl, Inches(6.8), Inches(2.4), Inches(6.35), Inches(4.8), accent=TEAL)
txb(sl, "SOCRATES Framework", Inches(6.95), Inches(2.5), Inches(6.0), Inches(0.35),
    size=13, color=MUTED, font="Consolas")
socrates = [("S","Site — Where exactly is the symptom?"),
            ("O","Onset — When did it start? Sudden or gradual?"),
            ("C","Character — Sharp, dull, burning, crushing?"),
            ("R","Radiation — Does it spread anywhere?"),
            ("A","Associated symptoms — Nausea, sweat, fever?"),
            ("T","Timing — Constant or intermittent?"),
            ("E","Exacerbating / relieving factors?"),
            ("S","Severity — Mild / Moderate / Severe")]
for i, (letter, desc) in enumerate(socrates):
    y = Inches(2.95 + i * 0.53)
    txb(sl, letter, Inches(6.95), y, Inches(0.35), Inches(0.4), size=16, bold=True, color=TEAL, font="Georgia")
    txb(sl, desc,   Inches(7.35), y, Inches(5.6),  Inches(0.4), size=13, color=MUTED, font="Calibri")

bottom_bar(sl)


# ── SLIDE 14: Safety & Critic Agent ────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Feature · PROJ-11")
slide_num(sl, 14)
title_line(sl, "Safety & ", "Critic Agent")
subtitle(sl, "Pre-LLM safety rules + full reasoning verification + hallucination detection")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

safety_items = [
    (ROSE,  "Class 1 — Cardiac Emergency",
     "Triggers: STEMI symptoms, crushing chest pain + radiation to jaw/arm, diaphoresis. Action: Immediate emergency escalation message before any AI reasoning."),
    (ROSE,  "Class 1 — Stroke Protocol",
     "Triggers: FAST criteria — facial droop, arm weakness, speech disturbance, sudden onset. Action: '999/911 immediately' prominent alert."),
    (AMBER, "Class 2 — Drug Interactions",
     "Triggers: anticoagulant + NSAID, MAOI + sympathomimetic, QT-prolonging combinations. Action: explicit interaction warning."),
    (TEAL,  "Class 3 — Reasoning Transparency",
     "Triggers: every diagnostic conclusion. Action: Critic agent reviews logic consistency, evidence alignment, hallucination markers."),
]
for i, (col, title, desc) in enumerate(safety_items):
    y = Inches(2.4 + i * 1.1)
    card_box(sl, Inches(0.5), y, Inches(7.8), Inches(1.0), accent=col)
    rect(sl, Inches(0.5), y, Inches(0.06), Inches(1.0), fill=col)
    txb(sl, title, Inches(0.7), y + Inches(0.1), Inches(7.4), Inches(0.35), size=15, bold=True, color=TEXT, font="Calibri")
    txb(sl, desc,  Inches(0.7), y + Inches(0.48), Inches(7.4), Inches(0.45), size=12, color=MUTED, wrap=True)

# Critic workflow (right side)
card_box(sl, Inches(8.5), Inches(2.4), Inches(4.65), Inches(4.45), accent=TEAL)
txb(sl, "Critic Agent Workflow", Inches(8.65), Inches(2.5), Inches(4.3), Inches(0.35),
    size=13, color=MUTED, font="Consolas")
workflow = ["① Receives: full Diagnostician output",
            "② Checks: logic consistency",
            "③ Checks: evidence alignment",
            "④ Checks: hallucination markers",
            "⑤ Checks: safety rule triggers",
            '⑥ Issues: "Consistent" or objection']
for i, step in enumerate(workflow):
    txb(sl, step, Inches(8.65), Inches(2.95 + i * 0.55),
        Inches(4.3), Inches(0.45), size=13, color=TEXT2)

card_box(sl, Inches(8.5), Inches(7.0), Inches(4.65), Inches(0.72), accent=ROSE)
txb(sl, "🛡  Pre-LLM Safety = O(1) · Pattern matching · Cannot be bypassed by prompt injection",
    Inches(8.65), Inches(7.08), Inches(4.3), Inches(0.55), size=11, color=ROSE, wrap=True)

bottom_bar(sl)


# ── SLIDE 15: Dual Auth Roles ──────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Feature · PROJ-9")
slide_num(sl, 15)
title_line(sl, "Dual Auth ", "Roles")
subtitle(sl, "JWT-based role-claim authentication — Patient and Provider with distinct access scopes")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

for i, (icon, role, col, perms, extra) in enumerate([
    ("🧑‍⚕️", "Patient Role", TEAL,
     ["Start new diagnostic sessions","Upload medical files (6 modalities)",
      "View own session history","Chat with AI Interviewer","Export own results PDF / JSON"],
     "Endpoints: /api/session/* · /api/patients/* · /api/analyze/*"),
    ("👨‍💼", "Provider Role", AMBER,
     ["Access ALL patient sessions","View cross-patient analytics",
      "Provider-only dashboard","Full audit trail with reasoning chains",
      "Patient management panel"],
     "Extra endpoints: /api/provider/sessions · /api/provider/analytics"),
]):
    x = Inches(0.5 + i * 6.4)
    card_box(sl, x, Inches(2.4), Inches(6.1), Inches(4.55), accent=col)
    txb(sl, icon, x + Inches(0.2), Inches(2.5), Inches(0.8), Inches(0.7), size=40)
    txb(sl, role, x + Inches(0.2), Inches(3.22), Inches(5.7), Inches(0.5),
        size=26, color=TEXT, font="Georgia")
    txb(sl, "JWT role claim · bcrypt hashed password",
        x + Inches(0.2), Inches(3.7), Inches(5.7), Inches(0.28), size=11, color=col, font="Consolas")
    hline(sl, x + Inches(0.2), Inches(4.02), Inches(5.7), col, Pt(0.5))
    for j, perm in enumerate(perms):
        txb(sl, f"→  {perm}", x + Inches(0.2), Inches(4.12 + j * 0.42),
            Inches(5.7), Inches(0.38), size=13, color=TEXT2)
    txb(sl, extra, x + Inches(0.2), Inches(6.2), Inches(5.7), Inches(0.55),
        size=11, color=col, font="Consolas", wrap=True)

card_box(sl, Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.6), accent=TEAL)
txb(sl, "Implementation: FastAPI OAuth2PasswordBearer · HS256-signed JWT · role field in payload · dependency injection get_current_user(role='provider') guards provider endpoints",
    Inches(0.65), Inches(7.15), Inches(12.0), Inches(0.45), size=11, color=MUTED, font="Consolas")

bottom_bar(sl)


# ── SLIDE 16: Multilingual Support ─────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Feature · PROJ-12")
slide_num(sl, 16)
title_line(sl, "Multilingual ", "Support")
subtitle(sl, "Full react-i18next internationalisation — UI, agent responses, voice transcription in EN / ZH / JA")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

for i, (flag, lang, code) in enumerate([("🇬🇧","English","en-US"),("🇨🇳","中文","zh-CN"),("🇯🇵","日本語","ja-JP")]):
    x = Inches(0.5 + i * 2.0)
    card_box(sl, x, Inches(2.4), Inches(1.8), Inches(1.4), accent=TEAL)
    txb(sl, flag, x + Inches(0.3), Inches(2.5), Inches(0.6), Inches(0.55), size=28)
    txb(sl, lang, x + Inches(0.1), Inches(3.0), Inches(1.6), Inches(0.35), size=16, color=TEXT, align=PP_ALIGN.CENTER)
    txb(sl, code, x + Inches(0.1), Inches(3.32), Inches(1.6), Inches(0.28), size=10, color=MUTED, align=PP_ALIGN.CENTER, font="Consolas")

# Coverage bars
card_box(sl, Inches(0.5), Inches(4.0), Inches(5.5), Inches(2.8), accent=TEAL)
txb(sl, "Translation Coverage", Inches(0.65), Inches(4.1), Inches(5.2), Inches(0.3), size=11, color=MUTED, font="Consolas")
coverages = [("UI Labels & Placeholders","100%",1.0),
             ("Error Messages & Disclaimers","100%",1.0),
             ("Body Parts & Duration Options","100%",1.0),
             ("Voice Transcription Language","3 langs",1.0)]
for i, (lbl, val, ratio) in enumerate(coverages):
    y = Inches(4.5 + i * 0.55)
    txb(sl, lbl, Inches(0.65), y, Inches(4.0), Inches(0.25), size=12, color=TEXT2)
    rect(sl, Inches(0.65), y + Inches(0.27), Inches(4.3 * ratio), Inches(0.2), fill=TEAL2)
    txb(sl, val,  Inches(4.97), y + Inches(0.25), Inches(0.8), Inches(0.22), size=11, color=TEAL, font="Consolas")

# Tech/UI/Audio steps
for i, (lbl, title, desc) in enumerate([
    ("Tech","react-i18next + i18next",
     "JSON locale files · useTranslation() hook · returnObjects for arrays · localStorage persist"),
    ("Audio","Speech Language Auto-Detection",
     "navigator.language → matched to supported lang code → passed to SpeechRecognition.lang"),
    ("UI","Seamless Switching",
     "Language toggle in navbar → full re-render with t() keys · no page reload"),
]):
    pipe_step(sl, lbl, title, desc, Inches(6.3), Inches(2.4 + i * 1.05), Inches(6.85))

stat_card(sl, "3",    "Supported languages", Inches(6.3), Inches(5.6), Inches(3.3), TEAL)
stat_card(sl, "200+", "Translated strings",  Inches(9.85), Inches(5.6), Inches(3.3), AMBER)

bottom_bar(sl)


# ── SLIDE 17: Session Management ───────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Feature · Data Persistence")
slide_num(sl, 17)
title_line(sl, "Session ", "Management")
subtitle(sl, "Automatic session persistence · full history access · one-click PDF/JSON export for clinical integration")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

db_steps = [
    ("SQLite","Four-Table Schema",
     "users(id, username, role, hashed_pw) · sessions(id, patient_id, created_at, symptoms) · messages(id, session_id, role, content) · uploads(id, session_id, filename, analysis)"),
    ("Auto","Zero-Effort Persistence",
     "Every session, message, and upload saved automatically · History page populates from DB on load · patient profile links maintained"),
    ("Export","PDF + JSON Clinical Export",
     "GET /api/session/{id}/export/pdf → formatted clinical report with diagnosis, reasoning, citations · /export/json → EHR integration"),
    ("Upload","File Association",
     "Files linked to session via session_id FK · analysis stored in DB · /api/sessions/{id}/uploads returns all files with analyses"),
]
for i, (num, title, desc) in enumerate(db_steps):
    pipe_step(sl, num, title, desc, Inches(0.5), Inches(2.4 + i * 1.02), Inches(7.5))

# Results tabs
card_box(sl, Inches(8.2), Inches(2.4), Inches(5.0), Inches(3.75), accent=TEAL)
txb(sl, "Results Page — 4 Tabs", Inches(8.35), Inches(2.5), Inches(4.7), Inches(0.3),
    size=13, color=MUTED, font="Consolas")
tabs = [("🩺","Diagnosis","Ranked differential · probability · Critic verdict", TEAL),
        ("🔗","Reasoning","Full agent chain · Interviewer → Diagnostician → Critic", AMBER),
        ("📖","Literature","PubMed citations · title · authors · year · relevance", ROSE),
        ("🖼","Media","Uploaded files · annotated images with overlay dots", MUTED)]
for i, (icon, tab, desc, col) in enumerate(tabs):
    y = Inches(2.9 + i * 0.8)
    rect(sl, Inches(8.35), y, Inches(0.08), Inches(0.6), fill=col)
    txb(sl, f"{icon}  {tab}", Inches(8.5), y + Inches(0.05), Inches(4.5), Inches(0.3),
        size=14, bold=True, color=col, font="Calibri")
    txb(sl, desc, Inches(8.5), y + Inches(0.35), Inches(4.5), Inches(0.38), size=11, color=MUTED, wrap=True)

stat_card(sl, "4", "DB tables",    Inches(8.2), Inches(6.3), Inches(2.35), TEAL)
stat_card(sl, "2", "Export formats",Inches(10.8), Inches(6.3), Inches(2.35), AMBER)

bottom_bar(sl)


# ── SLIDE 18: System Performance ───────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Evaluation")
slide_num(sl, 18)
title_line(sl, "System ", "Performance")
subtitle(sl, "Quantitative comparison across all major system dimensions — Demo A baseline vs current implementation")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

# Comparison table (left)
perf_rows = [
    ("Knowledge Corpus", "~5,000 docs", "47,441 docs (+848%)", TEAL),
    ("Embedding Model", "all-MiniLM 384d", "BioLORD-2023 768d", TEAL),
    ("Search Strategy", "Single cosine", "Dense+BM25+RRF", TEAL),
    ("Query Variants", "1 raw input", "3–4 LLM-rewritten", TEAL),
    ("Recall Threshold", "0.25 (strict)", "0.15 (+recall)", AMBER),
    ("Input Modalities", "1 (text)", "6 modalities", TEAL),
    ("Auth Roles", "None", "Patient + Provider", TEAL),
    ("UI Languages", "EN only", "EN / ZH / JA", TEAL),
    ("Image Annotation", "None", "9-region JSON overlay", TEAL),
    ("Safety Rules", "None", "3 classes", TEAL),
    ("Story Points", "13 SP", "48 SP (+35)", AMBER),
]
headers_p = ["Metric", "Demo A Baseline", "Current System"]
for ci, hdr in enumerate(headers_p):
    x = Inches(0.5 + ci * 2.55)
    txb(sl, hdr, x, Inches(2.38), Inches(2.5), Inches(0.35),
        size=11, bold=True, color=[TEXT2, ROSE, TEAL][ci], font="Consolas")
hline(sl, Inches(0.5), Inches(2.72), Inches(7.5), TEAL, Pt(0.5))
for ri, (metric, va, vb, col) in enumerate(perf_rows):
    y = Inches(2.8 + ri * 0.41)
    rect(sl, Inches(0.5), y, Inches(7.5), Inches(0.38),
         fill=RGBColor(0x0A,0x18,0x28) if ri%2==0 else BG2)
    txb(sl, metric, Inches(0.55), y + Inches(0.06), Inches(2.45), Inches(0.28), size=11, color=TEXT2)
    txb(sl, va, Inches(3.05), y + Inches(0.06), Inches(2.45), Inches(0.28), size=11, color=ROSE)
    txb(sl, vb, Inches(5.6),  y + Inches(0.06), Inches(2.45), Inches(0.28), size=11, color=col)

# Right: metric bars
card_box(sl, Inches(8.2), Inches(2.38), Inches(5.0), Inches(2.85), accent=TEAL)
txb(sl, "Key Improvement Ratios", Inches(8.35), Inches(2.48), Inches(4.7), Inches(0.3),
    size=11, color=MUTED, font="Consolas")
ratios = [("Knowledge Base Size","9.4× larger",0.94),
          ("Embedding Dims","2× denser",0.50),
          ("Query Coverage","3–4× more",0.80),
          ("Input Modalities","6× more",0.75),
          ("Story Points","+270%",0.90)]
for i, (lbl, val, ratio) in enumerate(ratios):
    y = Inches(2.88 + i * 0.45)
    txb(sl, lbl, Inches(8.35), y, Inches(3.2), Inches(0.25), size=11, color=TEXT2)
    rect(sl, Inches(8.35), y + Inches(0.27), Inches(4.5 * ratio), Inches(0.16), fill=TEAL2)
    txb(sl, val, Inches(8.35 + 4.5*ratio + 0.1), y + Inches(0.25), Inches(1.5), Inches(0.22),
        size=10, color=TEAL, font="Consolas")

# Sprint bar chart
card_box(sl, Inches(8.2), Inches(5.38), Inches(5.0), Inches(1.85), accent=AMBER)
txb(sl, "Scope Growth (Story Points)", Inches(8.35), Inches(5.48), Inches(4.7), Inches(0.28),
    size=11, color=MUTED, font="Consolas")
for i, (lbl, val, col, h) in enumerate([
    ("Sprint 1","13",ROSE,0.52),("Sprint 2","35",TEAL,1.0),("Total","48",AMBER,0.96)
]):
    x = Inches(9.0 + i * 1.4)
    rect(sl, x, Inches(7.23 - h * 0.9), Inches(0.9), Inches(h * 0.9), fill=col)
    txb(sl, val, x, Inches(7.23 - h * 0.9 - 0.28), Inches(0.9), Inches(0.28),
        size=13, color=col, align=PP_ALIGN.CENTER, font="Georgia")
    txb(sl, lbl, x, Inches(7.23), Inches(0.9), Inches(0.25), size=10, color=MUTED, align=PP_ALIGN.CENTER)

bottom_bar(sl)


# ── SLIDE 19: Technical Stack ──────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Implementation")
slide_num(sl, 19)
title_line(sl, "Technical ", "Stack")
subtitle(sl, "Technology choices justified by clinical AI requirements — performance, biomedical specificity, zero external dependencies")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

stacks = [
    ("Frontend", [
        ("⚛", "React + Vite", "Component-based UI · Fast HMR · Production bundle optimisation"),
        ("🌐", "react-i18next", "3-language support · JSON locale files · Runtime switching"),
        ("📷", "getUserMedia + Canvas", "Webcam frame capture · toBlob() JPEG · No extra library"),
        ("🎨", "Tailwind CSS + shadcn/ui", "Utility-first · Component library · Dark mode native"),
    ]),
    ("Backend", [
        ("⚡", "FastAPI + Uvicorn", "Async Python · Type-safe endpoints · Auto-generated OpenAPI docs"),
        ("🔐", "JWT + bcrypt", "HS256-signed tokens · Role claims · Password hashing"),
        ("🗄", "SQLite", "Zero-config persistence · 4 tables · FK relationships"),
        ("🩻", "pydicom", "DICOM file parsing · Pixel array extraction · HU normalisation"),
    ]),
    ("Knowledge / AI", [
        ("🔍", "Qdrant (embedded)", "Dense + Sparse hybrid · In-process · No network hop · RRF fusion"),
        ("🧬", "BioLORD-2023", "Biomedical ontology training · SNOMED CT · MeSH · UMLS · 768d"),
        ("📊", "BM25 Sparse Encoder", "Keyword matching · Complements dense vectors · Qdrant native"),
        ("🤖", "Anthropic claude-sonnet-4-6", "3-agent orchestration · Vision · Structured JSON output"),
    ]),
]
for ci, (col_title, items) in enumerate(stacks):
    x = Inches(0.5 + ci * 4.28)
    card_box(sl, x, Inches(2.4), Inches(4.1), Inches(4.75), accent=TEAL)
    txb(sl, col_title, x + Inches(0.15), Inches(2.5), Inches(3.8), Inches(0.35),
        size=13, color=TEAL, font="Consolas")
    hline(sl, x + Inches(0.15), Inches(2.87), Inches(3.8), TEAL, Pt(0.5))
    for i, (icon, name, desc) in enumerate(items):
        y = Inches(3.0 + i * 1.0)
        txb(sl, icon, x + Inches(0.15), y + Inches(0.05), Inches(0.4), Inches(0.42), size=20)
        txb(sl, name, x + Inches(0.6), y + Inches(0.0), Inches(3.4), Inches(0.35),
            size=14, bold=True, color=TEXT, font="Calibri")
        txb(sl, desc, x + Inches(0.6), y + Inches(0.38), Inches(3.4), Inches(0.48),
            size=11, color=MUTED, font="Calibri", wrap=True)

bottom_bar(sl)


# ── SLIDE 20: Design Principles ────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Engineering Quality")
slide_num(sl, 20)
title_line(sl, "Design ", "Principles")
subtitle(sl, "Architectural decisions that ensure clinical-grade reliability, safety, and extensibility")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

principles = [
    ("01","Safety-First — Pre-LLM Rules",
     "Hard-coded safety triggers evaluated before any LLM call. Pattern matching is O(1) and cannot be bypassed by prompt injection."),
    ("02","Grounded Outputs — No Fabrication",
     "Every diagnostic claim must be supported by a retrieved PubMed passage. Diagnostician explicitly instructed not to claim unsupported facts."),
    ("03","Agent Separation of Concerns",
     "Interviewer, Diagnostician, and Critic have distinct system prompts with no cross-contamination — minimises prompt injection surface."),
    ("04","Full Audit Trail",
     "Every agent's reasoning chain stored in SQLite. Provider can review, export, and dispute any AI output for clinical accountability."),
    ("05","Zero External Service Dependencies",
     "Qdrant embedded in-process, SQLite local, Web Speech API browser-native, Canvas API built-in — resilient, cost-free, offline-capable."),
    ("06","Progressive Disclosure",
     "Simple input for patients; full reasoning chain, citations, and provider analytics accessible without overwhelming the primary user flow."),
    ("07","Multilingual by Design",
     "i18n implemented from the start, not retrofitted. All strings externalised to JSON locale files — adding new languages is a single file."),
    ("08","Graceful Degradation",
     "All multimodal features degrade gracefully — if file analysis fails, session continues with text-only. No single point of failure."),
]

for i, (num, title, desc) in enumerate(principles):
    col_i = i % 2
    row_i = i // 2
    x = Inches(0.5 + col_i * 6.4)
    y = Inches(2.4 + row_i * 1.2)
    card_box(sl, x, y, Inches(6.1), Inches(1.12), accent=TEAL)
    txb(sl, num,   x + Inches(0.15), y + Inches(0.1), Inches(0.5), Inches(0.38),
        size=12, color=TEAL, font="Consolas")
    txb(sl, title, x + Inches(0.7),  y + Inches(0.08), Inches(5.2), Inches(0.35),
        size=15, bold=True, color=TEXT, font="Calibri")
    txb(sl, desc,  x + Inches(0.15), y + Inches(0.52), Inches(5.85), Inches(0.52),
        size=11, color=MUTED, wrap=True)

bottom_bar(sl)


# ── SLIDE 21: Limitations & Future Work ────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Honest Evaluation")
slide_num(sl, 21)
title_line(sl, "Limitations & ", "Future Work")
subtitle(sl, "Transparent assessment of current boundaries and roadmap for further research")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

lims = [
    (ROSE, "Evaluation Scope",
     "System evaluated on qualitative case studies and manual review. No standardised clinical NLP benchmark (MedQA, PubMedQA) run due to time constraints."),
    (ROSE, "Corpus Currency",
     "PubMed snapshot at ingestion time — new research published after ingest date not retrievable without re-running /api/rag/ingest."),
    (AMBER, "Latency",
     "3-agent sequential pipeline adds 8–15 s per diagnosis versus single-model approaches. Acceptable for non-urgent consultation; unsuitable for real-time triage."),
    (AMBER, "No EHR Integration",
     "System operates independently — no connection to hospital EMR/EHR systems. Clinical data must be entered manually."),
]
for i, (col, title, desc) in enumerate(lims):
    x = Inches(0.5 + (i % 2) * 6.4)
    y = Inches(2.4 + (i // 2) * 1.65)
    card_box(sl, x, y, Inches(6.1), Inches(1.55), accent=col)
    txb(sl, title, x + Inches(0.2), y + Inches(0.1), Inches(5.7), Inches(0.35),
        size=16, bold=True, color=TEXT, font="Calibri")
    txb(sl, desc, x + Inches(0.2), y + Inches(0.48), Inches(5.7), Inches(0.95),
        size=12, color=MUTED, wrap=True)

futures = [
    ("①", "Benchmark Evaluation", "Run against MedQA, PubMedQA, MIMIC-III to quantify accuracy gains over baseline LLM."),
    ("②", "On-Demand RAG Update", "Schedule PubMed delta ingest via /api/rag/ingest — keep knowledge currency without full re-ingest."),
    ("③", "Voice-First UI Rebuild", "Redesign around reliable audio pipeline — server-side Whisper transcription for production-grade voice input."),
    ("④", "HL7 FHIR Integration", "Export session data in FHIR R4 format for direct EHR/EMR import."),
]
txb(sl, "Future Directions", Inches(0.5), Inches(5.82), Inches(4), Inches(0.3),
    size=12, color=TEAL, font="Consolas")
for i, (num, title, desc) in enumerate(futures):
    x = Inches(0.5 + (i%2)*6.4)
    y = Inches(6.15 + (i//2)*0.55)
    txb(sl, f"{num} {title}:", x, y, Inches(3.0), Inches(0.3), size=13, bold=True, color=TEAL)
    txb(sl, desc, x + Inches(3.05), y, Inches(3.0), Inches(0.35), size=11, color=MUTED, wrap=True)

bottom_bar(sl)


# ── SLIDE 22: Key Contributions ────────────────────────────────────────────────
sl = add_slide()
eyebrow(sl, "Research Contribution")
slide_num(sl, 22)
title_line(sl, "Key ", "Contributions")
subtitle(sl, "What this project demonstrates that existing solutions do not")
hline(sl, Inches(0.6), Inches(2.15), Inches(3))

contribs = [
    (TEAL,  "Multi-Agent Clinical Simulation",
     "First demonstration of a three-specialist agent pipeline (Interviewer → Diagnostician → Critic) specifically designed to mirror multi-disciplinary clinical team decision-making with safety verification.",
     "Addresses: hallucination, single-point reasoning failure, missing safety layer"),
    (AMBER, "Multimodal-Enriched RAG",
     "Novel pipeline combining LLM-rewritten medical queries (patient language → SNOMED/MeSH), BioLORD-2023 biomedical embeddings, BM25 sparse search, and RRF fusion on a 47,441-document PubMed corpus. Image findings extracted via LLM and injected into RAG queries.",
     "Addresses: knowledge currency, retrieval precision, multimodal gap"),
    (ROSE,  "Empathetic Clinical Dialogue",
     "Systematic application of SOCRATES framework within an LLM-powered agent with explicit empathy conventions — warm language, concern acknowledgement before clinical probing — alongside multilingual support.",
     "Addresses: patient-centred design, accessibility, language barrier"),
    (MUTED, "Explainable Safety Architecture",
     "Full reasoning audit trail per session stored in SQLite, accessible via Results → Reasoning tab. Pre-LLM pattern-matched safety rules guarantee O(1) escalation independent of model behaviour.",
     "Addresses: clinical accountability, auditability, regulatory compliance"),
]

for i, (col, title, desc, addr) in enumerate(contribs):
    col_i = i % 2
    row_i = i // 2
    x = Inches(0.5 + col_i * 6.4)
    y = Inches(2.4 + row_i * 2.4)
    card_box(sl, x, y, Inches(6.1), Inches(2.28), accent=col)
    rect(sl, x, y, Inches(0.08), Inches(2.28), fill=col)
    txb(sl, title, x + Inches(0.2), y + Inches(0.1), Inches(5.7), Inches(0.38),
        size=17, bold=True, color=TEXT, font="Calibri")
    txb(sl, desc, x + Inches(0.2), y + Inches(0.5), Inches(5.7), Inches(1.1),
        size=12, color=MUTED, wrap=True)
    txb(sl, addr, x + Inches(0.2), y + Inches(1.65), Inches(5.7), Inches(0.5),
        size=10, color=col, font="Consolas", wrap=True)

bottom_bar(sl)


# ── SLIDE 23: Thank You ────────────────────────────────────────────────────────
sl = add_slide()
txb(sl, "COMP9900 · Team 9900-W18C-CAKE · UNSW Sydney",
    Inches(0.6), Inches(0.55), Inches(12), Inches(0.35),
    size=12, color=TEAL, font="Consolas", align=PP_ALIGN.CENTER)

# Main title
txb(sl, "Medi", Inches(2.5), Inches(1.5), Inches(3.5), Inches(2.0), size=100, color=TEXT, font="Georgia")
txb(sl, "Chain", Inches(5.35), Inches(1.5), Inches(5.0), Inches(2.0), size=100, italic=True, color=TEAL, font="Georgia")

txb(sl, "A Multi-Agent Collaborative System for Empathetic Clinical Diagnostic Reasoning",
    Inches(1.0), Inches(3.5), Inches(11), Inches(0.6),
    size=19, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

hline(sl, Inches(2.0), Inches(4.25), Inches(9), TEAL)

# Stats grid
final_stats = [
    ("47,441","Medical Documents in RAG",TEAL),
    ("48 SP","Total Story Points",TEAL),
    ("6","Input Modalities",AMBER),
    ("3","UI Languages",AMBER),
    ("9.4×","Knowledge Base Growth",ROSE),
    ("768d","BioLORD Embedding Dims",MUTED),
]
for i, (val, lbl, col) in enumerate(final_stats):
    x = Inches(0.5 + (i%3) * 4.28)
    y = Inches(4.6 + (i//3) * 1.2)
    card_box(sl, x, y, Inches(4.1), Inches(1.1), accent=col)
    txb(sl, val, x + Inches(0.15), y + Inches(0.05), Inches(3.8), Inches(0.6),
        size=36, color=col, font="Georgia")
    txb(sl, lbl, x + Inches(0.15), y + Inches(0.65), Inches(3.8), Inches(0.38),
        size=11, color=MUTED, font="Consolas")

txb(sl, "Thank you",
    Inches(0.5), Inches(7.1), Inches(12), Inches(0.35),
    size=14, color=RGBColor(0x30,0x60,0x80), italic=True,
    align=PP_ALIGN.CENTER, font="Georgia")

bottom_bar(sl)


# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "MediChain_DemoB.pptx")
prs.save(out)
print(f"✓ Saved: {out}")
print(f"  Slides: {len(prs.slides)}")
